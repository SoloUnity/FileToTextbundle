
import os
from pptx import Presentation
from io import BytesIO
import re
from PIL import Image

def pptx_to_textbundle(pptx_path):
    # Determine the output directory name
    output_dir_name = os.path.splitext(os.path.basename(pptx_path))[0] + ".textbundle"
    output_dir_path = os.path.join(os.path.dirname(pptx_path), output_dir_name)
    
    # Create directories for TextBundle and assets
    assets_dir = os.path.join(output_dir_path, "assets")
    os.makedirs(assets_dir, exist_ok=True)
    
    # Extract content from PowerPoint
    presentation = Presentation(pptx_path)
    
    # Extract title from the first slide to use as filename
    first_slide_title = "text"
    if presentation.slides and presentation.slides[0].shapes.title:
        first_slide_title = presentation.slides[0].shapes.title.text
    markdown_filename = f"{first_slide_title}.md"
    
    markdown_content = []
    for slide_num, slide in enumerate(presentation.slides):
        slide_text = " ".join([shape.text for shape in slide.shapes if shape.has_text_frame])
        
        # Skip slide if it contains the word "poll"
        if "poll" in slide_text.lower():
            continue
        
        # Default to using "Slide [number]" as subtitle
        subtitle = f"Slide {slide_num + 1}"
        
        # Check if slide has a title and use it as subtitle if present
        if slide.shapes.title:
            subtitle = slide.shapes.title.text
        
        markdown_content.append(f"## {subtitle}\n")
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:  # Avoid duplicating title
                # Convert point form and numbered lists to markdown format
                paragraphs = shape.text_frame.paragraphs
                for paragraph in paragraphs:
                    prefix = ""
                    text = paragraph.text.strip().replace("•", "").strip()
                    # Check if the text starts with a number followed by a dot (indicating a numbered list)
                    if re.match(r"^\d+\.", text):
                        prefix = text.split(".")[0] + ". "
                        text = text.split(".")[1].strip()
                    elif paragraph.level > 0 or "•" in paragraph.text:  # Bullet point
                        prefix = "  " * paragraph.level + "- "
                    markdown_content.append(prefix + text)
                markdown_content.append("\n")
            elif shape.shape_type == 13:  # Shape type 13 corresponds to Picture
                image_stream = BytesIO(shape.image.blob)
                image_format = shape.image.ext
                image = Image.open(image_stream)
                if image.width < 100 or image.height < 100:  # Skip images smaller than 100x100 pixels
                    continue
                image_name = f"slide_{slide_num + 1}_image.{image_format.lower()}"
                image_path = os.path.join(assets_dir, image_name)
                with open(image_path, "wb") as image_file:
                    image_file.write(image_stream.getvalue())
                # Add image reference to the markdown content
                markdown_content.append(f"![{subtitle} Image](assets/{image_name})\n\n")
    
    # Save markdown content
    with open(os.path.join(output_dir_path, markdown_filename), "w") as f:
        f.write("\n".join(markdown_content))
    
    print(f"TextBundle created at {output_dir_path}")

if __name__ == "__main__":
    import sys
    if len(sys.argv) != 2:
        print("Usage: python pptx_to_textbundle_with_actual_image_size_filter.py [path_to_pptx]")
        sys.exit(1)
    pptx_path = sys.argv[1]
    pptx_to_textbundle(pptx_path)

